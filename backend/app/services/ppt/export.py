import os
from typing import List
import logging

import io
import img2pdf
from pptx import Presentation
from pptx.util import Inches

logger = logging.getLogger(__name__)

class ExportService:

    @staticmethod
    def create_pptx_from_images(
        image_paths: List[str],
        output_path: str = None,
    ) -> bytes:

        if not isinstance(image_paths, list) or len(image_paths) == 0:
            logger.error("图像路径为空或者非列表")
            raise ValueError("图像路径必须为非空列表")

        prs = Presentation()

        # 比例16:9
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

        # 添加图片
        added_slides = 0
        for image_path in image_paths:
            if not isinstance(image_path, str):
                logger.warning(f"无效的图像路径: {image_path}")
                continue

            if not os.path.exists(image_path):
                logger.warning(f"图像不存在: {image_path}")
                continue

            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.add_picture(
                image_path,
                Inches(0),
                Inches(0),
                Inches(10),
                Inches(5.625),
            )
            added_slides += 1

            logger.info(f"已添加 {added_slides} 张图片")

            if output_path:
                prs.save(output_path)
                logger.info(f"已保存到 {output_path}")
                return None
            else:
                pptx_bytes = io.BytesIO()
                prs.save(pptx_bytes)
                pptx_bytes.seek(0)
                return pptx_bytes.getvalue()


    def create_pdf_from_images(
        image_paths: List[str],
        output_path: str = None,
    ) -> bytes:
        if not isinstance(image_paths, list) or len(image_paths) == 0:
            logger.error("图像路径为空或者非列表")
            raise ValueError("图像路径必须为非空列表")

        layout_fun = img2pdf.get_layout_fun(16, 9)
        pdf_bytes = img2pdf.convert(image_paths, layout_fun=layout_fun)
        if output_path:
            with open(output_path, "wb") as f:
                f.write(pdf_bytes)
            logger.info(f"已保存到 {output_path}")
            return None
        else:
            return pdf_bytes


                


            


